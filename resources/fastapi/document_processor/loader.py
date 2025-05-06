import mimetypes
import threading
from typing import List, Tuple
from langchain_community.document_loaders import PyMuPDFLoader, TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import CharacterTextSplitter, MarkdownTextSplitter
from pptx import Presentation
from docx import Document as DocxDocument
import fitz
import aiohttp
import asyncio
import base64
import os
import subprocess


from document_processor.file_to_markdown import generate_markdown_from_file
from image_processor.utils import get_page_images_from_pdf_async

PDF_MIME_TYPES = ["application/pdf"]
TEXT_MIME_TYPES = ["text/plain"]
POWERPOINT_TYPES = [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
]
WORD_TYPES = [
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
]
SPREADSHEET_TYPES = ["text/csv", "application/csv"]
UPLOAD_ACCEPTED_DOCUMENTS = (
    PDF_MIME_TYPES + TEXT_MIME_TYPES + POWERPOINT_TYPES + WORD_TYPES + SPREADSHEET_TYPES
)
DECOMPOSE_ACCEPTED_DOCUMENTS = (
    PDF_MIME_TYPES + TEXT_MIME_TYPES
)

class DocumentsLoader:

    def __init__(self, documents: List[str]):
        self._document_paths = documents

        self._documents: List[Document] = []
        self._splitted_documents: List[Document] = []
        self._images: List[List[str]] = []

        self._markdown_splitter = MarkdownTextSplitter(chunk_size=500, chunk_overlap=50)
        self._text_splitter = CharacterTextSplitter(
            separator="/n", chunk_size=500, chunk_overlap=50
        )

    @property
    def documents(self):
        return self._documents

    @property
    def splitted_documents(self):
        return self._splitted_documents

    @property
    def images(self):
        return self._images

    async def  load_documents(
        self,
        temp_dir: str,
        split_documents: bool = False,
        load_markdown: bool = True,
        load_images: bool = False,
        presigned_file_urls: List[str] = None,
    ):
        documents: List[Document] = []
        images: List[str] = []

        splitted_documents: List[Document] = []
        for file_path, presigned_file_url in zip(self._document_paths, presigned_file_urls):
            mime_type = mimetypes.guess_type(file_path)[0]
            if mime_type not in DECOMPOSE_ACCEPTED_DOCUMENTS:
                continue

            docs = []
            imgs = []
            if mime_type in PDF_MIME_TYPES:
                docs, imgs = await self.load_pdf(
                    file_path, load_markdown, load_images, temp_dir, presigned_file_url
                )

            elif mime_type in TEXT_MIME_TYPES:
                docs = self.load_text(file_path)

            documents.extend(docs)
            images.append(imgs)

            if split_documents:
                splitted_documents.extend(self.split_documents(docs, mime_type))

        self._documents = documents
        self._splitted_documents = splitted_documents
        self._images = images

    

    def split_documents(self, documents: List[Document], mime_type):
        # if mime_type in PDF_MIME_TYPES:
        #     return self._markdown_splitter.split_documents(documents)
        return self._text_splitter.split_documents(documents)
 
    def clip_longer_documents(self, documents: List[Document], clip_after: int = 1200):
        for document in documents:
            document.page_content = document.page_content[:clip_after]
        return documents

    async def load_pdf(
        self, file_path: str, load_markdown: bool, load_images: bool, temp_dir: str, presigned_file_url: str
    ) -> Tuple[List[Document], List[str]]:
        image_paths = []
        documents: List[Document] = []

        if load_markdown:
            doc = fitz.open(file_path)
            total_pages = len(doc)
            metadata = {
                "source": file_path,
                "total_pages": total_pages,
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "subject": doc.metadata.get("subject", ""),
            }
            # if total_pages > 10:
            import time
            t1 = time.time()
            full_text = await self.decompose_pdf_to_markdown(total_pages, presigned_file_url)
            print(f"Time Taken for pdf to md: {time.time() - t1}")
            documents.append(
                Document(page_content=full_text.strip(), metadata=metadata)
            )
            # else:
            #     markdown = generate_markdown_from_file(file_path)
            #     documents.append(
            #         Document(page_content=markdown.markdown, metadata=metadata)
            #     )

        if load_images:
            image_paths = await get_page_images_from_pdf_async(file_path, temp_dir)

        return documents, image_paths

    async def decompose_pdf_to_markdown(self,total_pages: int, source_url: str) -> str:
        async def fetch_markdown(session, pages, headers):
            params = {'source': source_url, 'pages': ','.join(map(str, pages))}
            async with session.get('https://surajbeston--pdf-to-md-pymupdfllm-decompose.modal.run', params=params, headers=headers) as response:
                data = await response.json()
                return data['data']

        async with aiohttp.ClientSession() as session:
            proxy_auth = base64.b64encode(f"{os.getenv('MODAL_ID')}:{os.getenv('MODAL_SECRET')}".encode()).decode()
            headers = {"Proxy-Authorization": f"Basic {proxy_auth}"}
            tasks = []
            batch_size = 5
            for start_page in range(0, total_pages, batch_size):
                end_page = min(start_page + batch_size, total_pages)
                pages = list(range(start_page, end_page))
                tasks.append(fetch_markdown(session, pages, headers))

            results = await asyncio.gather(*tasks)

        return '\n'.join(results)
    

    def load_text(self, file_path: str) -> List[Document]:
        loader = TextLoader(file_path)
        return loader.load()

    def load_msword(self, file_path: str) -> List[Document]:
        document = DocxDocument(file_path)
        text = "\n".join([paragraph.text for paragraph in document.paragraphs])
        return [Document(page_content=text)]

    def load_powerpoint(self, file_path: str) -> List[Document]:
        presentation = Presentation(file_path)

        extracted_text = ""
        for index, slide in enumerate(presentation.slides):
            extracted_text += f"# Slide {index + 1}\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        extracted_text += f"{paragraph.text}\n"
                    extracted_text += "\n"
            extracted_text += "\n\n"
        return [Document(page_content=extracted_text)]
